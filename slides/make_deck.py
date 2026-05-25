#!/usr/bin/env python3
"""
VoyagerAI — Auth0 for AI Agents demo deck.
Audience: CISOs + security architects. ~20 slides, English, Auth0 brand.

Run:
    pip install python-pptx
    python3 make_deck.py
    open voyagerai-deck.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ─── Auth0 brand ───────────────────────────────────────────────────────────────
BG_DARK    = RGBColor(0x0F, 0x14, 0x19)   # near-black, slight cool tint
BG_PANEL   = RGBColor(0x16, 0x1C, 0x24)   # one shade lighter for panels
BG_CODE    = RGBColor(0x0A, 0x0F, 0x14)
TEXT       = RGBColor(0xF4, 0xF4, 0xF4)
TEXT_DIM   = RGBColor(0xA0, 0xA9, 0xB4)
TEXT_MUTED = RGBColor(0x6B, 0x73, 0x7E)
BORDER     = RGBColor(0x2A, 0x32, 0x3C)
ORANGE     = RGBColor(0xEB, 0x54, 0x24)   # Auth0 brand
ORANGE_DIM = RGBColor(0x99, 0x37, 0x18)
SUCCESS    = RGBColor(0x4A, 0xC4, 0x86)
WARN       = RGBColor(0xE8, 0xA8, 0x3A)
DANGER     = RGBColor(0xE5, 0x5A, 0x5A)
INFO       = RGBColor(0x6F, 0xA8, 0xDC)

# ─── Layout constants ──────────────────────────────────────────────────────────
SLIDE_W = Inches(13.333)   # 16:9 widescreen
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.55)
TITLE_TOP = Inches(0.45)
TITLE_H = Inches(0.7)

FONT_FAMILY = "Calibri"
MONO_FAMILY = "Consolas"

# ─── Helpers ───────────────────────────────────────────────────────────────────

def add_blank_slide(prs):
    """Add a blank slide with the dark background."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_DARK
    bg.line.fill.background()
    bg.shadow.inherit = False
    return slide

def add_text(slide, left, top, width, height, text, *, size=18, color=TEXT,
             bold=False, italic=False, font=FONT_FAMILY, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, line_spacing=1.2):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    return tb

def add_rich_paragraphs(slide, left, top, width, height, paragraphs, *,
                        font=FONT_FAMILY, anchor=MSO_ANCHOR.TOP):
    """paragraphs: list of (text, size, color, bold, italic, line_spacing)."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    for i, spec in enumerate(paragraphs):
        text, size, color, bold, italic, line_spacing = spec
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = text
        r.font.name = font
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.bold = bold
        r.font.italic = italic
    return tb

def add_panel(slide, left, top, width, height, *, fill=BG_PANEL, border=BORDER, accent=None):
    """Rounded panel with optional left accent bar."""
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    panel.adjustments[0] = 0.04
    panel.fill.solid()
    panel.fill.fore_color.rgb = fill
    panel.line.color.rgb = border
    panel.line.width = Pt(0.75)
    panel.shadow.inherit = False
    if accent:
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.06), height)
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent
        bar.line.fill.background()
        bar.shadow.inherit = False
    return panel

def add_chip(slide, left, top, text, *, color=ORANGE, fg=None):
    """Pill-shaped chip with text."""
    width = Inches(0.05 + 0.085 * len(text))
    height = Inches(0.28)
    chip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    chip.adjustments[0] = 0.5
    chip.fill.solid()
    chip.fill.fore_color.rgb = color
    chip.line.fill.background()
    chip.shadow.inherit = False
    tf = chip.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.name = FONT_FAMILY
    r.font.size = Pt(10)
    r.font.bold = True
    r.font.color.rgb = fg if fg else BG_DARK
    return chip, width

def add_title(slide, text, subtitle=None):
    """Standard slide title strip with subtle orange underline."""
    add_text(slide, MARGIN, TITLE_TOP, SLIDE_W - 2*MARGIN, TITLE_H, text,
             size=30, color=TEXT, bold=True)
    underline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                        MARGIN, TITLE_TOP + Inches(0.78),
                                        Inches(0.6), Inches(0.04))
    underline.fill.solid(); underline.fill.fore_color.rgb = ORANGE
    underline.line.fill.background(); underline.shadow.inherit = False
    if subtitle:
        add_text(slide, MARGIN, TITLE_TOP + Inches(0.95),
                 SLIDE_W - 2*MARGIN, Inches(0.4),
                 subtitle, size=14, color=TEXT_DIM, italic=True)

def add_footer(slide, page_num, total):
    add_text(slide, MARGIN, SLIDE_H - Inches(0.4),
             Inches(4), Inches(0.3),
             "VoyagerAI · Auth0 for AI Agents", size=9, color=TEXT_MUTED)
    add_text(slide, SLIDE_W - MARGIN - Inches(1.5), SLIDE_H - Inches(0.4),
             Inches(1.5), Inches(0.3),
             f"{page_num} / {total}", size=9, color=TEXT_MUTED, align=PP_ALIGN.RIGHT)

def add_code_block(slide, left, top, width, height, lines):
    panel = add_panel(slide, left, top, width, height, fill=BG_CODE)
    tb = slide.shapes.add_textbox(left + Inches(0.18), top + Inches(0.12),
                                   width - Inches(0.36), height - Inches(0.24))
    tf = tb.text_frame
    tf.word_wrap = False
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.18
        r = p.add_run()
        r.text = line if line else " "
        r.font.name = MONO_FAMILY
        r.font.size = Pt(11)
        r.font.color.rgb = TEXT_DIM
    return panel

def add_bullets(slide, left, top, width, height, items, *, size=14, gap=0.18):
    """items: list of strings or (label, body) tuples."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.35
        p.space_after = Pt(int(gap * 72))
        # bullet dot
        rd = p.add_run(); rd.text = "·  "
        rd.font.name = FONT_FAMILY; rd.font.size = Pt(size + 4)
        rd.font.color.rgb = ORANGE; rd.font.bold = True
        if isinstance(item, tuple):
            label, body = item
            rl = p.add_run(); rl.text = label
            rl.font.name = FONT_FAMILY; rl.font.size = Pt(size)
            rl.font.color.rgb = TEXT; rl.font.bold = True
            rb = p.add_run(); rb.text = "  " + body
            rb.font.name = FONT_FAMILY; rb.font.size = Pt(size)
            rb.font.color.rgb = TEXT_DIM
        else:
            rb = p.add_run(); rb.text = item
            rb.font.name = FONT_FAMILY; rb.font.size = Pt(size)
            rb.font.color.rgb = TEXT
    return tb

# ─── Slide builders ────────────────────────────────────────────────────────────

def slide_title(prs):
    s = add_blank_slide(prs)
    # Big orange band at top
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.15))
    band.fill.solid(); band.fill.fore_color.rgb = ORANGE
    band.line.fill.background(); band.shadow.inherit = False
    # eyebrow
    add_text(s, MARGIN, Inches(2.2), Inches(8), Inches(0.4),
             "AUTH0 FOR AI AGENTS  ·  LIVE DEMO", size=14, color=ORANGE, bold=True)
    # title
    add_text(s, MARGIN, Inches(2.7), Inches(11), Inches(1.5),
             "VoyagerAI", size=64, color=TEXT, bold=True)
    add_text(s, MARGIN, Inches(3.7), Inches(11.5), Inches(0.8),
             "Identity, authorization, and accountability for the agentic enterprise.",
             size=22, color=TEXT_DIM)
    # subtle credit chip bottom
    add_text(s, MARGIN, Inches(6.7), Inches(8), Inches(0.3),
             "For CISOs and security architects  ·  ~25 minutes  ·  Three acts, fully live",
             size=11, color=TEXT_MUTED)
    return s

def slide_problem(prs):
    s = add_blank_slide(prs)
    add_title(s, "The problem AI agents create")
    add_text(s, MARGIN, Inches(1.7), Inches(12), Inches(0.6),
             "Agents need credentials. Today's options all fail at scale.",
             size=18, color=TEXT_DIM, italic=True)
    cards = [
        ("Static API keys", "One leaked key, one compromised agent — and a global blast radius. No expiry, no scope, no audit trail per call.", DANGER),
        ("Shared OAuth scopes", "Bolt 'agent_admin' onto a user's access token and you've created an over-privileged account that bypasses every policy you wrote for humans.", WARN),
        ("Custom mTLS / VPN", "Solves transport, not authorization. The agent still acts with someone else's identity, and you still can't answer 'who did what'.", WARN),
    ]
    y = Inches(2.4)
    for label, body, accent in cards:
        add_panel(s, MARGIN, y, Inches(12.2), Inches(1.25), accent=accent)
        add_text(s, MARGIN + Inches(0.3), y + Inches(0.2),
                 Inches(11.6), Inches(0.4),
                 label, size=18, color=TEXT, bold=True)
        add_text(s, MARGIN + Inches(0.3), y + Inches(0.6),
                 Inches(11.6), Inches(0.7),
                 body, size=13, color=TEXT_DIM)
        y += Inches(1.4)
    return s

def slide_thesis(prs):
    s = add_blank_slide(prs)
    add_title(s, "What this demo shows",
              "Three CISO-friendly stories. One platform — Auth0 — wires them all.")
    bullets = [
        ("The agent acts as me, with limits.",
         "First-party agent + step-up MFA + Manager CIBA + bounded-authority cap."),
        ("The agent only sees what I see.",
         "Third-party agent with consent + Auth0 FGA enforcing relationship-based authorization."),
        ("The agent never holds my password.",
         "Token Vault brokers a short-lived, scoped token at the moment of need."),
    ]
    y = Inches(2.0)
    for label, body in bullets:
        add_panel(s, MARGIN, y, Inches(12.2), Inches(1.4), accent=ORANGE)
        add_text(s, MARGIN + Inches(0.35), y + Inches(0.2),
                 Inches(11.6), Inches(0.5),
                 label, size=20, color=TEXT, bold=True)
        add_text(s, MARGIN + Inches(0.35), y + Inches(0.7),
                 Inches(11.6), Inches(0.7),
                 body, size=14, color=TEXT_DIM)
        y += Inches(1.55)
    return s

def slide_architecture(prs):
    s = add_blank_slide(prs)
    add_title(s, "Architecture at a glance",
              "Two services for VoyagerAI; one tenant for Auth0; one mock downstream.")
    # Three columns: Portal | API | Auth0
    box_w = Inches(3.7); box_h = Inches(2.2)
    gutter = Inches(0.18)
    y = Inches(2.0)
    boxes = [
        ("voyagerai-portal", "React + Vite · :3000",
         ["• Travel Agent (1st-party)", "• Personal AI Assistant (3rd-party)",
          "• Token Inspector · Audit · MCP timeline"]),
        ("voyagerai-api",    "Express · :3002",
         ["• OAuth 2 Resource Server (RS256)", "• Agent loop · 6 MCP-style tools",
          "• Audit log · SSE event bus"]),
        ("Auth0 tenant",     "OIDC · CIBA · FGA · Vault",
         ["• 1 API + 1 SPA + 2 M2M apps", "• Action: amr + max_trip_value",
          "• FGA store · Mgmt API"]),
    ]
    x = MARGIN
    for title, sub, lines in boxes:
        add_panel(s, x, y, box_w, box_h, accent=ORANGE)
        add_text(s, x + Inches(0.25), y + Inches(0.2),
                 box_w - Inches(0.4), Inches(0.4),
                 title, size=16, color=TEXT, bold=True)
        add_text(s, x + Inches(0.25), y + Inches(0.6),
                 box_w - Inches(0.4), Inches(0.3),
                 sub, size=11, color=ORANGE, italic=True)
        add_bullets(s, x + Inches(0.25), y + Inches(1.0),
                    box_w - Inches(0.4), Inches(1.1), lines, size=11)
        x += box_w + gutter
    # Below: "downstream services" row
    add_text(s, MARGIN, Inches(4.5), Inches(12), Inches(0.4),
             "Downstream / federated", size=12, color=TEXT_MUTED, bold=True)
    down_y = Inches(4.9); down_h = Inches(1.5)
    items = [
        ("VoyagerVault",   "Mock SaaS · separate audience\nReceives short-lived tokens via\nclient_credentials + private_key_jwt.", ORANGE),
        ("Manager device",  "Auth0 Guardian app\nReceives binding_message via CIBA\n(RFC 9126 / OIDC backchannel).", INFO),
        ("Mgmt API",       "Auth0 Management API\nGrants list/revoke for the\nConnected Agents page.", SUCCESS),
    ]
    x = MARGIN
    for label, body, accent in items:
        add_panel(s, x, down_y, box_w, down_h, accent=accent)
        add_text(s, x + Inches(0.25), down_y + Inches(0.18),
                 box_w - Inches(0.4), Inches(0.4),
                 label, size=14, color=TEXT, bold=True)
        add_text(s, x + Inches(0.25), down_y + Inches(0.55),
                 box_w - Inches(0.4), Inches(0.9),
                 body, size=11, color=TEXT_DIM)
        x += box_w + gutter
    return s

def slide_act1_overview(prs):
    s = add_blank_slide(prs)
    add_chip(s, MARGIN, Inches(0.55), "ACT  1", color=ORANGE)
    add_title(s, "“The agent acts as me, with limits.”",
              "First-party Travel Agent · 3-tier authorization · CIBA · bounded authority.")
    # Tier matrix
    rows = [
        ("Tier 1",  "≤ $500",         "Instant",                  "Has book:travel scope. Logged.",                   SUCCESS),
        ("Tier 2",  "$500 – $2,000",  "Step-up MFA",              "Token must carry amr=['mfa']. Mirrored from ID by Auth0 Action.", WARN),
        ("Tier 3",  "> $2,000",       "Manager CIBA approval",    "Phone vibrates. Manager approves with biometrics in Auth0 Guardian.", INFO),
        ("Tier 3+", "> max_trip_value","Bounded authority blocks","Even a manager-approved booking is denied if it exceeds the user's per-trip cap (custom claim).", DANGER),
    ]
    y = Inches(2.1)
    for tier, amt, gate, body, accent in rows:
        add_panel(s, MARGIN, y, Inches(12.2), Inches(1.0), accent=accent)
        add_text(s, MARGIN + Inches(0.25), y + Inches(0.18),
                 Inches(1.6), Inches(0.4),
                 tier, size=14, color=accent, bold=True)
        add_text(s, MARGIN + Inches(1.9), y + Inches(0.18),
                 Inches(2.0), Inches(0.4),
                 amt, size=14, color=TEXT, bold=True)
        add_text(s, MARGIN + Inches(4.0), y + Inches(0.18),
                 Inches(2.6), Inches(0.4),
                 gate, size=14, color=TEXT, bold=True)
        add_text(s, MARGIN + Inches(0.25), y + Inches(0.58),
                 Inches(11.7), Inches(0.4),
                 body, size=11, color=TEXT_DIM)
        y += Inches(1.13)
    return s

def slide_act1_ciba(prs):
    s = add_blank_slide(prs)
    add_chip(s, MARGIN, Inches(0.55), "ACT  1  ·  DETAIL", color=ORANGE)
    add_title(s, "CIBA, not push — the agent never asks for the manager's password",
              "Auth0's Client-Initiated Backchannel Authentication (RFC 9126).")
    # Two-column: flow steps + sequence diagram
    add_text(s, MARGIN, Inches(2.0), Inches(6.3), Inches(0.4),
             "What happens", size=14, color=ORANGE, bold=True)
    flow = [
        ("1.", "Agent calls /bc-authorize with login_hint={iss,sub of manager}, binding_message, scope=approve:travel."),
        ("2.", "Auth0 sends a Guardian push to the manager's phone."),
        ("3.", "Manager approves with biometrics. No browser, no password reuse."),
        ("4.", "Agent polls /oauth/token with grant_type=urn:openid:params:grant-type:ciba."),
        ("5.", "Auth0 issues a manager-bound access token. Booking finalizes."),
    ]
    y = Inches(2.5)
    for n, body in flow:
        add_text(s, MARGIN, y, Inches(0.4), Inches(0.4),
                 n, size=14, color=ORANGE, bold=True)
        add_text(s, MARGIN + Inches(0.4), y, Inches(5.9), Inches(0.7),
                 body, size=12, color=TEXT_DIM, line_spacing=1.3)
        y += Inches(0.8)
    # Right side: code excerpt
    add_text(s, Inches(7.2), Inches(2.0), Inches(5.7), Inches(0.4),
             "lib/ciba.js · 1st-party agent", size=14, color=ORANGE, bold=True)
    code = [
        "const body = new URLSearchParams({",
        "  client_id: AUTH0_AGENT_CLIENT_ID,",
        "  client_assertion_type:",
        "    'urn:ietf:params:oauth:client-assertion-type:jwt-bearer',",
        "  client_assertion: buildClientAssertion(...),",
        "  scope: 'openid approve:travel',",
        "  binding_message: sanitize(",
        "    'Approve flight to Tokyo for USD 4,500'),",
        "  login_hint: JSON.stringify({",
        "    format: 'iss_sub',",
        "    iss: ISSUER,",
        "    sub: AUTH0_MANAGER_USER_ID,",
        "  }),",
        "});",
    ]
    add_code_block(s, Inches(7.2), Inches(2.5), Inches(5.7), Inches(4.4), code)
    return s

def slide_act1_bounded(prs):
    s = add_blank_slide(prs)
    add_chip(s, MARGIN, Inches(0.55), "ACT  1  ·  DETAIL", color=ORANGE)
    add_title(s, "Bounded authority — even the manager can't override the org cap",
              "A custom JWT claim, set at login, enforced server-side.")
    add_text(s, MARGIN, Inches(2.0), Inches(6.3), Inches(0.4),
             "Login Action injects the cap", size=14, color=ORANGE, bold=True)
    code1 = [
        "// Auth0 Action · Login Flow",
        "exports.onExecutePostLogin = async (event, api) => {",
        "  // Cap stays under the org's policy. Per-user later",
        "  // by reading event.user.app_metadata.max_trip_value.",
        "  api.accessToken.setCustomClaim(",
        "    'https://voyagerai.demo/max_trip_value',",
        "    5000,",
        "  );",
        "  // mirror ID-token amr into the access token so",
        "  // server-side step-up checks see it.",
        "  api.accessToken.setCustomClaim(",
        "    'https://voyagerai.demo/amr',",
        "    event.authentication?.methods?.map(m => m.name) ?? [],",
        "  );",
        "};",
    ]
    add_code_block(s, MARGIN, Inches(2.5), Inches(6.3), Inches(4.4), code1)
    add_text(s, Inches(7.2), Inches(2.0), Inches(5.7), Inches(0.4),
             "API enforces it after the manager approves", size=14, color=ORANGE, bold=True)
    code2 = [
        "// lib/policy.js  — runs on EVERY booking",
        "if (cibaApproved && amount > maxTripValue) {",
        "  return {",
        "    decision: 'bounded_authority_exceeded',",
        "    bounded_authority: { max: maxTripValue,",
        "                         requested: amount },",
        "  };",
        "}",
        "",
        "// CISO punchline:",
        "// The manager OK'd it. The org's cap still wins.",
        "// One claim, audited, signed, in every token.",
    ]
    add_code_block(s, Inches(7.2), Inches(2.5), Inches(5.7), Inches(4.4), code2)
    return s

def slide_act2(prs):
    s = add_blank_slide(prs)
    add_chip(s, MARGIN, Inches(0.55), "ACT  2", color=ORANGE)
    add_title(s, "“The agent only sees what I see.”",
              "3rd-party app + consent + Auth0 FGA (relationship-based access control).")
    # Two-column
    add_text(s, MARGIN, Inches(2.0), Inches(6.3), Inches(0.4),
             "OAuth scopes answer ‘can the agent call the tool?’", size=14, color=ORANGE, bold=True)
    add_bullets(s, MARGIN, Inches(2.5), Inches(6.3), Inches(2.2), [
        ("Personal AI Assistant", "Separate Auth0 app. Allowed scopes: read:trips, read:expenses. NOT book:travel."),
        ("Consent screen", "Auth0 prompts the user once with the exact scopes before issuing a token."),
        ("Hard boundary", "An attempt to call book_travel returns insufficient_scope at the API. The LLM cannot escalate."),
    ], size=12)
    add_text(s, Inches(7.2), Inches(2.0), Inches(5.7), Inches(0.4),
             "FGA answers ‘on which records?’", size=14, color=ORANGE, bold=True)
    code = [
        "model",
        "  schema 1.1",
        "",
        "type user",
        "type cost_center",
        "  relations",
        "    define member: [user]",
        "type user_profile",
        "  relations",
        "    define owner: [user]",
        "    define cost_center: [cost_center]",
        "    define can_view: owner or member from cost_center",
        "",
        "// → Lara (eng)  → can_view ✓",
        "// → VP Eng (exec) → can_view ✗",
    ]
    add_code_block(s, Inches(7.2), Inches(2.5), Inches(5.7), Inches(4.4), code)
    return s

def slide_act3(prs):
    s = add_blank_slide(prs)
    add_chip(s, MARGIN, Inches(0.55), "ACT  3", color=ORANGE)
    add_title(s, "“The agent never holds my password.”",
              "Token Vault — Auth0 brokers a short-lived, audience-scoped token at the moment of need.")
    # Sequence list left, payload right
    add_text(s, MARGIN, Inches(2.0), Inches(6.3), Inches(0.4),
             "Flow", size=14, color=ORANGE, bold=True)
    flow = [
        ("a.", "User asks the agent to save trip notes to VoyagerVault."),
        ("b.", "Agent calls Auth0 /oauth/token  — client_credentials + private_key_jwt + audience=https://api.voyagervault.demo."),
        ("c.", "Auth0 issues a vault-scoped token (write:vault). 1-hour TTL. Distinct audience from main API."),
        ("d.", "Agent calls VoyagerVault with Authorization: Bearer <vault token>  AND  X-On-Behalf-Of: <user JWT>."),
        ("e.", "Vault validates BOTH tokens. Stores entry keyed by user sub."),
    ]
    y = Inches(2.5)
    for n, body in flow:
        add_text(s, MARGIN, y, Inches(0.4), Inches(0.4),
                 n, size=13, color=ORANGE, bold=True)
        add_text(s, MARGIN + Inches(0.4), y, Inches(5.9), Inches(0.6),
                 body, size=11, color=TEXT_DIM, line_spacing=1.3)
        y += Inches(0.65)
    # Right column
    add_text(s, Inches(7.2), Inches(2.0), Inches(5.7), Inches(0.4),
             "Payload — what the agent sends", size=14, color=ORANGE, bold=True)
    code = [
        "POST /oauth/token",
        "  grant_type=client_credentials",
        "  audience=https://api.voyagervault.demo",
        "  client_id=<travel-agent>",
        "  client_assertion_type=jwt-bearer",
        "  client_assertion=<RS256 JWT signed",
        "                    with private key>",
        "  scope=write:vault read:vault",
        "",
        "→ access_token  (TTL 3600s)",
        "  aud: api.voyagervault.demo",
        "  sub: <agent-client-id>",
        "  scope: write:vault",
        "",
        "// CISO punchline:",
        "// No static credential. No reuse across",
        "// services. Every token audited.",
    ]
    add_code_block(s, Inches(7.2), Inches(2.5), Inches(5.7), Inches(4.4), code)
    return s

def slide_features_matrix(prs):
    s = add_blank_slide(prs)
    add_title(s, "Auth0 surface area used",
              "Ten capabilities, one tenant, ~5 hours of integration work.")
    cells = [
        ("OIDC + PKCE",         "User login. Standard SPA + Auth0Provider."),
        ("RBAC",                "Roles 'Traveler' and 'Manager'. 'approve:travel' gated by role."),
        ("Custom claims",       "Action injects max_trip_value + amr mirror."),
        ("private_key_jwt",     "1st-party Travel Agent M2M auth (RFC 7521)."),
        ("Step-up MFA",         "acr_values=mfa during Tier 2 retries."),
        ("CIBA (RFC 9126)",     "Manager approval via Guardian for Tier 3."),
        ("FGA",                 "Relationship-based read of user_profile records."),
        ("Token Vault",         "Audience-scoped downstream token (mock)."),
        ("Mgmt API",            "Connected Agents listing + grant revocation."),
        ("Refresh Token Rotation","Long-lived sessions without long-lived tokens."),
    ]
    cols = 2
    cell_w = Inches(6.0); cell_h = Inches(0.9)
    x0 = MARGIN; y0 = Inches(2.0)
    gutter_x = Inches(0.2); gutter_y = Inches(0.15)
    for i, (label, body) in enumerate(cells):
        col = i % cols
        row = i // cols
        x = x0 + (cell_w + gutter_x) * col
        y = y0 + (cell_h + gutter_y) * row
        add_panel(s, x, y, cell_w, cell_h, accent=ORANGE)
        add_text(s, x + Inches(0.2), y + Inches(0.13),
                 cell_w - Inches(0.3), Inches(0.4),
                 label, size=14, color=TEXT, bold=True)
        add_text(s, x + Inches(0.2), y + Inches(0.48),
                 cell_w - Inches(0.3), Inches(0.4),
                 body, size=11, color=TEXT_DIM)
    return s

def slide_agent_identity(prs):
    s = add_blank_slide(prs)
    add_chip(s, MARGIN, Inches(0.55), "BUILD  ·  AGENT IDENTITY", color=ORANGE)
    add_title(s, "Three agents, three identity stories",
              "Each agent gets its own Auth0 application. No shared secrets.")
    rows = [
        ("Travel Agent (1st-party · M2M)",
         "private_key_jwt (RFC 7521). RSA key pair generated locally; public JWKS pasted in Auth0. Agent signs a fresh client_assertion JWT for every /oauth/token call. Auth0 verifies the signature against the JWKS — no client_secret in env.",
         SUCCESS),
        ("Personal AI Assistant (3rd-party · SPA)",
         "Separate Auth0 application requiring user consent. Limited to read:trips, read:expenses. Portal uses a second @auth0/auth0-spa-js client (isolated cache) to keep tokens disjoint from the 1st-party agent.",
         INFO),
        ("Mgmt M2M (system)",
         "client_credentials with classic client_secret. Used only by the Express API for read:grants, delete:grants, read:clients on the Auth0 Management API. Never user-facing.",
         WARN),
    ]
    y = Inches(2.0)
    for label, body, accent in rows:
        add_panel(s, MARGIN, y, Inches(12.2), Inches(1.5), accent=accent)
        add_text(s, MARGIN + Inches(0.3), y + Inches(0.2),
                 Inches(11.6), Inches(0.4),
                 label, size=15, color=TEXT, bold=True)
        add_text(s, MARGIN + Inches(0.3), y + Inches(0.6),
                 Inches(11.6), Inches(0.85),
                 body, size=11, color=TEXT_DIM, line_spacing=1.35)
        y += Inches(1.65)
    return s

def slide_policy_code(prs):
    s = add_blank_slide(prs)
    add_chip(s, MARGIN, Inches(0.55), "BUILD  ·  POLICY", color=ORANGE)
    add_title(s, "The 3-tier policy is a pure function",
              "lib/policy.js — testable, reusable in serverless, no side effects.")
    code = [
        "function evaluate({ amountUSD, jwtPayload, cibaApproved = false }) {",
        "  const amr = jwtPayload['https://voyagerai.demo/amr'] || jwtPayload.amr || [];",
        "  const maxTripValue = jwtPayload[BOUNDED_AUTHORITY_CLAIM] ?? Infinity;",
        "",
        "  // Tier 1 — instant",
        "  if (amountUSD <= 500) return { decision: 'allow', tier: 1 };",
        "",
        "  // Tier 2 — step-up MFA",
        "  if (amountUSD <= 2000) {",
        "    if (amr.includes('mfa')) return { decision: 'allow', tier: 2 };",
        "    return { decision: 'requires_stepup', tier: 2 };",
        "  }",
        "",
        "  // Tier 3 — manager CIBA + bounded authority",
        "  if (!cibaApproved) return { decision: 'requires_ciba', tier: 3 };",
        "  if (amountUSD > maxTripValue) {",
        "    return {",
        "      decision: 'bounded_authority_exceeded',",
        "      bounded_authority: { max: maxTripValue, requested: amountUSD },",
        "      tier: 3,",
        "    };",
        "  }",
        "  return { decision: 'allow', tier: 3, bounded_authority: { max: maxTripValue } };",
        "}",
    ]
    add_code_block(s, MARGIN, Inches(2.0), SLIDE_W - 2*MARGIN, Inches(4.9), code)
    return s

def slide_fga_detail(prs):
    s = add_blank_slide(prs)
    add_chip(s, MARGIN, Inches(0.55), "BUILD  ·  FGA", color=ORANGE)
    add_title(s, "FGA — relationship-based access control",
              "Google Zanzibar-style policy. The agent inherits the user's read perimeter.")
    # Left: model
    add_text(s, MARGIN, Inches(2.0), Inches(6.3), Inches(0.4),
             "Model", size=14, color=ORANGE, bold=True)
    code1 = [
        "type user",
        "type cost_center",
        "  relations",
        "    define member: [user]",
        "type user_profile",
        "  relations",
        "    define owner: [user]",
        "    define cost_center: [cost_center]",
        "    define can_view:",
        "      owner or member from cost_center",
    ]
    add_code_block(s, MARGIN, Inches(2.5), Inches(6.3), Inches(2.6), code1)
    # Tuples
    add_text(s, MARGIN, Inches(5.3), Inches(6.3), Inches(0.4),
             "Demo tuples", size=14, color=ORANGE, bold=True)
    code2 = [
        "user:traveler  member  cost_center:engineering",
        "cost_center:engineering  cost_center  user_profile:peer-eng",
        "cost_center:executive    cost_center  user_profile:vp-engineering",
    ]
    add_code_block(s, MARGIN, Inches(5.8), Inches(6.3), Inches(1.2), code2)
    # Right: outcomes
    add_text(s, Inches(7.2), Inches(2.0), Inches(5.7), Inches(0.4),
             "Outcomes", size=14, color=ORANGE, bold=True)
    outcomes = [
        ("user:traveler  can_view  user_profile:traveler",  "ALLOW · owner",                         SUCCESS),
        ("user:traveler  can_view  user_profile:peer-eng",  "ALLOW · member from cost_center",        SUCCESS),
        ("user:traveler  can_view  user_profile:vp-engineering", "DENY · different cost_center",     DANGER),
    ]
    y = Inches(2.5)
    for q, ans, accent in outcomes:
        add_panel(s, Inches(7.2), y, Inches(5.7), Inches(0.95), accent=accent)
        add_text(s, Inches(7.4), y + Inches(0.13),
                 Inches(5.4), Inches(0.4),
                 q, size=10, color=TEXT, font=MONO_FAMILY)
        add_text(s, Inches(7.4), y + Inches(0.5),
                 Inches(5.4), Inches(0.4),
                 ans, size=12, color=accent, bold=True)
        y += Inches(1.1)
    add_text(s, Inches(7.2), Inches(6.3), Inches(5.7), Inches(0.6),
             "FGA decisions stream live into the MCP timeline — auditors see the same yes/no the agent saw.",
             size=11, color=TEXT_MUTED, italic=True)
    return s

def slide_token_vault_detail(prs):
    s = add_blank_slide(prs)
    add_chip(s, MARGIN, Inches(0.55), "BUILD  ·  TOKEN VAULT", color=ORANGE)
    add_title(s, "Token Vault — service-agnostic by design",
              "We mock VoyagerVault. The same code targets Google Calendar / Slack / Salesforce.")
    # Two-column comparison
    cell_w = Inches(6.0); cell_h = Inches(4.6)
    add_panel(s, MARGIN, Inches(2.0), cell_w, cell_h, accent=ORANGE)
    add_panel(s, MARGIN + cell_w + Inches(0.2), Inches(2.0), cell_w, cell_h, accent=SUCCESS)
    add_text(s, MARGIN + Inches(0.25), Inches(2.15),
             cell_w - Inches(0.4), Inches(0.4),
             "Demo (mock)", size=14, color=ORANGE, bold=True)
    add_text(s, MARGIN + Inches(0.25), Inches(2.55),
             cell_w - Inches(0.4), Inches(0.4),
             "VoyagerVault — separate-audience M2M", size=13, color=TEXT, bold=True)
    code1 = [
        "POST /oauth/token",
        "  grant_type=client_credentials",
        "  audience=https://api.voyagervault.demo",
        "  scope=write:vault read:vault",
        "  client_assertion=<RS256 JWT>",
        "",
        "// returns aud=voyagervault, scope=write:vault",
    ]
    add_code_block(s, MARGIN + Inches(0.25), Inches(3.0),
                   cell_w - Inches(0.5), Inches(2.1), code1)
    add_text(s, MARGIN + Inches(0.25), Inches(5.3),
             cell_w - Inches(0.4), Inches(1.0),
             "Captures: audience scoping, JWT signing, scope enforcement, audit. Doesn't capture: a third-party token-store.",
             size=11, color=TEXT_DIM, line_spacing=1.4)
    # Right (production)
    pr_x = MARGIN + cell_w + Inches(0.2)
    add_text(s, pr_x + Inches(0.25), Inches(2.15),
             cell_w - Inches(0.4), Inches(0.4),
             "Production (real Token Vault)", size=14, color=SUCCESS, bold=True)
    add_text(s, pr_x + Inches(0.25), Inches(2.55),
             cell_w - Inches(0.4), Inches(0.4),
             "Federated grant + connection token", size=13, color=TEXT, bold=True)
    code2 = [
        "POST /oauth/token",
        "  grant_type=urn:auth0:params:oauth:grant-type:",
        "             token-exchange:federated-connection-",
        "             access-token",
        "  subject_token=<user JWT>",
        "  connection=google-oauth2",
        "  client_assertion=<RS256 JWT>",
        "",
        "// returns Google access token from vault",
    ]
    add_code_block(s, pr_x + Inches(0.25), Inches(3.0),
                   cell_w - Inches(0.5), Inches(2.1), code2)
    add_text(s, pr_x + Inches(0.25), Inches(5.3),
             cell_w - Inches(0.4), Inches(1.0),
             "Same call site — only the grant string and the destination URL change. The agent never reads a Google credential.",
             size=11, color=TEXT_DIM, line_spacing=1.4)
    return s

def slide_observability(prs):
    s = add_blank_slide(prs)
    add_chip(s, MARGIN, Inches(0.55), "BUILD  ·  OBSERVABILITY", color=ORANGE)
    add_title(s, "Observable surface — what auditors see",
              "Three real-time channels, one append-only audit log.")
    panels = [
        ("MCP Event Timeline",
         "SSE-fed live feed of every OAuth call (token request, /bc-authorize, FGA check, vault exchange). Color-coded by kind.",
         "/api/mcp/events"),
        ("Token Inspector",
         "Side-by-side decoded JWTs for every agent (1st-party, 3rd-party). Annotated claims. Scopes, custom claims, amr, audience.",
         "client-side jwt-decode"),
        ("Audit Trail",
         "Every tool call logged with agent_sub, user_sub, tool, decision, scopes, FGA result, CIBA id, bounded_authority.",
         "audit.log.jsonl + /api/audit"),
        ("Connected Agents",
         "Live Mgmt API call listing every 3rd-party app the user authorized. One-click revoke (DELETE /grants/:id).",
         "/api/v2/grants"),
    ]
    cell_w = Inches(6.0); cell_h = Inches(2.0)
    gutter = Inches(0.2)
    for i, (label, body, source) in enumerate(panels):
        col = i % 2; row = i // 2
        x = MARGIN + (cell_w + gutter) * col
        y = Inches(2.0) + (cell_h + Inches(0.18)) * row
        add_panel(s, x, y, cell_w, cell_h, accent=ORANGE)
        add_text(s, x + Inches(0.25), y + Inches(0.18),
                 cell_w - Inches(0.4), Inches(0.4),
                 label, size=15, color=TEXT, bold=True)
        add_text(s, x + Inches(0.25), y + Inches(0.6),
                 cell_w - Inches(0.4), Inches(1.0),
                 body, size=11, color=TEXT_DIM, line_spacing=1.4)
        add_text(s, x + Inches(0.25), y + Inches(1.55),
                 cell_w - Inches(0.4), Inches(0.35),
                 source, size=10, color=ORANGE, font=MONO_FAMILY, italic=True)
    return s

def slide_principles(prs):
    s = add_blank_slide(prs)
    add_title(s, "The principles, in one slide")
    bullets = [
        ("Agent identity ≠ user identity",
         "Every agent is its own Auth0 application. Acting 'on behalf of' the user is a token claim, not a credential share."),
        ("Scopes gate the tool. FGA gates the record.",
         "OAuth answers 'can the agent call book_travel?'. FGA answers 'on which user_profile?'. They compose; you need both."),
        ("Brokered tokens, not stored credentials",
         "Auth0 mints a fresh, audience-scoped, time-bounded token at the moment of need. No long-lived secrets in agent code or env."),
        ("Audit everything, in real time",
         "Every authorization event — accept and reject — flows to a timeline the auditor and the security team can subscribe to today."),
    ]
    y = Inches(2.0)
    for label, body in bullets:
        add_panel(s, MARGIN, y, Inches(12.2), Inches(1.15), accent=ORANGE)
        add_text(s, MARGIN + Inches(0.3), y + Inches(0.18),
                 Inches(11.6), Inches(0.4),
                 label, size=15, color=TEXT, bold=True)
        add_text(s, MARGIN + Inches(0.3), y + Inches(0.55),
                 Inches(11.6), Inches(0.6),
                 body, size=12, color=TEXT_DIM, line_spacing=1.4)
        y += Inches(1.27)
    return s

def slide_what_next(prs):
    s = add_blank_slide(prs)
    add_title(s, "What's next — and what we deferred",
              "Phased delivery; the high-impact items are live, the rest are scoped.")
    # Two-column
    add_text(s, MARGIN, Inches(2.0), Inches(6.3), Inches(0.4),
             "Live in the demo today", size=14, color=SUCCESS, bold=True)
    add_bullets(s, MARGIN, Inches(2.5), Inches(6.3), Inches(4.4), [
        "OIDC + PKCE login (Traveler · Manager)",
        "1st-party Travel Agent (M2M · private_key_jwt)",
        "3-tier policy + step-up MFA + bounded authority",
        "True CIBA (RFC 9126) for Tier 3",
        "3rd-party Personal AI Assistant + consent",
        "Auth0 FGA (relationship-based authz)",
        "Token Vault (mock VoyagerVault)",
        "Connected Agents (Mgmt API list + revoke)",
        "Token Inspector · Audit Trail · MCP timeline",
    ], size=12)
    add_text(s, Inches(7.2), Inches(2.0), Inches(5.7), Inches(0.4),
             "Deferred · easy to layer on", size=14, color=ORANGE, bold=True)
    add_bullets(s, Inches(7.2), Inches(2.5), Inches(5.7), Inches(4.4), [
        ("Settings · Agent Control Panel",
         "Per-scope toggles, bounded-authority slider, revoke-all."),
        ("Time-bounded 3rd-party agent",
         "time_bound claim + middleware that rejects outside window."),
        ("Tool Authorization screen",
         "‘Missing required scopes’ UI when Settings removes a scope."),
        ("Real Google Calendar via Token Vault",
         "Same code path as the mock; just swap the grant string."),
    ], size=12)
    return s

def slide_qa(prs):
    s = add_blank_slide(prs)
    # Centered Q&A
    add_text(s, Inches(0), Inches(2.5), SLIDE_W, Inches(1.5),
             "Questions?", size=72, color=TEXT, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0), Inches(4.0), SLIDE_W, Inches(0.6),
             "Try it live, ask anything about the wiring.",
             size=18, color=TEXT_DIM, italic=True, align=PP_ALIGN.CENTER)
    underline = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(6.4), Inches(3.85),
                                    Inches(0.6), Inches(0.06))
    underline.fill.solid(); underline.fill.fore_color.rgb = ORANGE
    underline.line.fill.background(); underline.shadow.inherit = False
    return s

# ─── Compose ───────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    builders = [
        slide_title,            # 1
        slide_problem,          # 2
        slide_thesis,           # 3
        slide_architecture,     # 4
        slide_act1_overview,    # 5
        slide_act1_ciba,        # 6
        slide_act1_bounded,     # 7
        slide_act2,             # 8
        slide_act3,             # 9
        slide_features_matrix,  # 10
        slide_agent_identity,   # 11
        slide_policy_code,      # 12
        slide_fga_detail,       # 13
        slide_token_vault_detail, # 14
        slide_observability,    # 15
        slide_principles,       # 16
        slide_what_next,        # 17
        slide_qa,               # 18
    ]
    total = len(builders)
    for i, build in enumerate(builders, start=1):
        slide = build(prs)
        if i > 1:  # skip footer on title
            add_footer(slide, i, total)

    out = "voyagerai-deck.pptx"
    prs.save(out)
    print(f"Wrote {out}  ({total} slides)")

if __name__ == "__main__":
    main()
